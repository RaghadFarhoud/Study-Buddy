from __future__ import annotations

import os
import uuid
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from app.domain.models import (
    CanonicalDocument,
    Section,
    SourceType,
    BlockType,
    HeadingBlock,
    ParagraphBlock,
    BulletListBlock,
    NumberedListBlock,
    FigureBlock,
    TableBlock,
    TableCell,
    ListItem,
)
class PptxExtractor:
    def __init__(self, image_output_dir: str = "output/images") -> None:
        self.image_output_dir = Path(image_output_dir)
        self.image_output_dir.mkdir(parents=True, exist_ok=True)

    def extract(self, file_path: str) -> CanonicalDocument:
        prs = Presentation(file_path)

        document_id = Path(file_path).stem
        sections: List[Section] = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            section = self._extract_slide_section(
                slide=slide,
                slide_index=slide_index,
            )
            sections.append(section)

        return CanonicalDocument(
            document_id=document_id,
            source_type=SourceType.PPTX,
            source_file=os.path.basename(file_path),
            title=document_id,
            language=None,
            sections=sections,
            metadata={
                "extractor": "PptxExtractor",
                "slide_count": len(prs.slides),
            },
        )
    def _expand_paragraph_lines(self, cleaned_paragraphs):
        expanded = []

        for para, text in cleaned_paragraphs:
            parts = [part.strip() for part in text.splitlines() if part.strip()]
            if parts:
                for part in parts:
                    expanded.append((para, part))

        return expanded
    def _extract_slide_section(self, slide, slide_index: int) -> Section:
        raw_items = []

        for shape in slide.shapes:
            item = self._extract_shape(shape=shape, slide_index=slide_index)
            if item is not None:
                raw_items.append(item)

        raw_items.sort(key=lambda item: (item["top"], item["left"]))

        blocks = []
        detected_title: Optional[str] = None
        
        for order, item in enumerate(raw_items, start=1):
            block = item["block"]
            block.order = order

            if block.type == BlockType.HEADING and detected_title is None:
                if getattr(block, "text", None):
                    detected_title = block.text.strip()

            blocks.append(block)

        if  detected_title is None and len(blocks) == 1:
            only_block = blocks[0]
            block_text = getattr(only_block, "text", None)
            if block_text and len(block_text.split()) <= 6 and len(block_text) <= 80:
                detected_title = block_text.strip()

        section_title = detected_title or f"Slide {slide_index}"

        for block in blocks:
            block.section_path = [section_title]

        notes_text = self._extract_slide_notes(slide)

        return Section(
            title=section_title,
            path=[section_title],
            blocks=blocks,
            page_or_slide=slide_index,
            notes=notes_text,
            metadata={
                "slide_index": slide_index,
                "shape_count": len(slide.shapes),
            },
        )

    def _extract_shape(self, shape: BaseShape, slide_index: int):

        if self._should_skip_shape(shape):
            return None 
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._extract_picture(shape, slide_index)

        if getattr(shape, "has_table", False):
            return self._extract_table(shape, slide_index)

        if getattr(shape, "has_text_frame", False):
            return self._extract_text_shape(shape, slide_index)

        return None

    def _extract_text_shape(self, shape: BaseShape, slide_index: int):
        paragraphs = getattr(shape.text_frame, "paragraphs", [])
        cleaned_paragraphs = []

        # expanded_paragraphs = self._expand_paragraph_lines(cleaned_paragraphs)

        for para in paragraphs:
            text = para.text.strip()
            if text:
                cleaned_paragraphs.append((para, text))

        if not cleaned_paragraphs:
            fallback_text = getattr(shape, "text", "")
            fallback_text = fallback_text.strip() if fallback_text else ""
            if fallback_text:
                cleaned_paragraphs.append((None, fallback_text))
        expanded_paragraphs = self._expand_paragraph_lines(cleaned_paragraphs)
        bbox = {
            "left": int(shape.left),
            "top": int(shape.top),
            "width": int(shape.width),
            "height": int(shape.height),
        }

        common_metadata = {
            "shape_name": getattr(shape, "name", None),
            "shape_type": str(shape.shape_type),
            "is_placeholder": getattr(shape, "is_placeholder", False),
            "bbox": bbox,
        }

        # 1) heading
        if self._is_title_shape(shape, cleaned_paragraphs):
            text = "\n".join(text for _, text in expanded_paragraphs)
            block = HeadingBlock(
                id=uuid.uuid4().hex,
                order=0,
                text=text,
                level=1,
                page_or_slide=slide_index,
                section_path=[],
                metadata=common_metadata,
            )
            return {"top": int(shape.top), "left": int(shape.left), "block": block}

        # 2) numbered list
        if self._is_numbered_list(expanded_paragraphs):
            items = [
                ListItem(text=text, level=getattr(para, "level", 0) or 0)
                for para, text in expanded_paragraphs
            ]
            
            shape_name = (getattr(shape, "name", "") or "").lower()
            if "slide number" in shape_name:
                return None


            block = NumberedListBlock(
                id=uuid.uuid4().hex,
                order=0,
                items=items,
                page_or_slide=slide_index,
                section_path=[],
                metadata=common_metadata,
            )
            return {"top": int(shape.top), "left": int(shape.left), "block": block}

        # 3) bullet list
        if self._is_bullet_list(expanded_paragraphs):
            items = [
                ListItem(
                    text=self._strip_bullet_prefix(text),
                    level=getattr(para, "level", 0) or 0,
                )
                for para, text in expanded_paragraphs
            ]
            block = BulletListBlock(
                id=uuid.uuid4().hex,
                order=0,
                items=items,
                page_or_slide=slide_index,
                section_path=[],
                metadata=common_metadata,
            )
            return {"top": int(shape.top), "left": int(shape.left), "block": block}

        # 4) paragraph
        merged_text = "\n".join(text for _, text in cleaned_paragraphs).strip()

        if not merged_text:
            return None
        block = ParagraphBlock(
            id=uuid.uuid4().hex,
            order=0,
            text=merged_text,
            page_or_slide=slide_index,
            section_path=[],
            metadata=common_metadata,
        )
        return {"top": int(shape.top), "left": int(shape.left), "block": block}
    def _should_skip_shape(self, shape: BaseShape) -> bool:

        shape_name = (getattr(shape, "name", "") or "").lower()

        if "slide number" in shape_name:
            return True

        if "date placeholder" in shape_name:
            return True

        if "footer placeholder" in shape_name:
            return True

        return False
    def _extract_picture(self, shape: BaseShape, slide_index: int):
        image = shape.image
        ext = image.ext
        filename = f"slide_{slide_index}_{uuid.uuid4().hex[:8]}.{ext}"
        image_path = self.image_output_dir / filename

        with open(image_path, "wb") as f:
            f.write(image.blob)

        block = FigureBlock(
            id=uuid.uuid4().hex,
            order=0,
            page_or_slide=slide_index,
            section_path=[],
            text=None,
            image_path=str(image_path),
            image_name=filename,
            caption=None,
            rel_id=None,
            alt_text=None,
            metadata={
                "shape_name": getattr(shape, "name", None),
                "shape_type": str(shape.shape_type),
                "bbox": {
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                },
            },
        )

        return {
            "top": int(shape.top),
            "left": int(shape.left),
            "block": block,
        }

    def _extract_table(self, shape: BaseShape, slide_index: int):
        table = shape.table
        cells: List[TableCell] = []

        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                cells.append(
                    TableCell(
                        row=r_idx,
                        col=c_idx,
                        text=cell.text.strip(),
                    )
                )

        rows = len(table.rows)
        cols = len(table.columns)

        block = TableBlock(
            id=uuid.uuid4().hex,
            order=0,
            page_or_slide=slide_index,
            section_path=[],
            text=None,
            rows=rows,
            cols=cols,
            cells=cells,
            summary=f"جدول يحتوي على {rows} صفوف و {cols} أعمدة.",
            metadata={
                "shape_name": getattr(shape, "name", None),
                "shape_type": str(shape.shape_type),
                "bbox": {
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                },
            },
        )

        return {
            "top": int(shape.top),
            "left": int(shape.left),
            "block": block,
        }

    def _is_title_shape(self, shape: BaseShape, cleaned_paragraphs) -> bool:
        if getattr(shape, "is_placeholder", False):
            placeholder_format = getattr(shape, "placeholder_format", None)
            if placeholder_format is not None:
                placeholder_type = str(placeholder_format.type).lower()
                if "title" in placeholder_type:
                    return True

        if len(cleaned_paragraphs) == 1:
            text = cleaned_paragraphs[0][1]
            words = text.split()
            if len(words) <= 6 and len(text) <= 60 and int(shape.top) < 1200000:
                return True

        return False

    def _is_numbered_list(self, cleaned_paragraphs) -> bool:
        if not cleaned_paragraphs:
            return False

        matches = 0
        for _, text in cleaned_paragraphs:
            if self._looks_like_numbered_item(text):
                matches += 1

        return matches >= 1 and len(cleaned_paragraphs) >= 2

    def _is_bullet_list(self, cleaned_paragraphs) -> bool:
        if not cleaned_paragraphs or len(cleaned_paragraphs) < 2:
            return False

        bullet_like = 0
        short_lines = 0
        for para, text in cleaned_paragraphs:
            stripped = text.strip()
            if stripped.startswith(("•", "-", "▪", "◦", "*")):
                bullet_like += 1

            if len(stripped.split()) <= 12:
                short_lines += 1
        if bullet_like >=1:
            return True
        if short_lines ==len(cleaned_paragraphs):
            return True

        return False

    def _strip_bullet_prefix(self, text: str) -> str:
        prefixes = ("•", "-", "▪", "◦", "*")
        stripped = text.strip()
        for prefix in prefixes:
            if stripped.startswith(prefix):
                return stripped[len(prefix):].strip()
        return stripped

    def _looks_like_numbered_item(self, text: str) -> bool:
        import re
        return bool(re.match(r"^\d+[\.\)]\s*", text.strip()))

    def _extract_slide_notes(self, slide) -> Optional[str]:
        try:
            notes_slide = slide.notes_slide
        except Exception:
            return None

        if not notes_slide:
            return None

        texts = []
        for shape in notes_slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)

        if not texts:
            return None

        merged = "\n".join(texts).strip()
        lines = [line.strip() for line in merged.splitlines() if line.strip()]
        while lines and lines[-1].isdigit():
            lines.pop()
        if not lines:
            return None
        merged = "\n".join(lines).strip()

        if not merged or merged.isdigit():
            return None
        

        return merged